import io
import uuid
import json
import pandas as pd
import matplotlib
matplotlib.use("Agg")  # headless — no display backend available in a serverless function
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

import storage


# ---------------- LOAD FILE (in-memory, no disk) ---------------- #
def load_file_from_upload(file_storage):
    """file_storage is a werkzeug FileStorage from request.files — read
    directly from the upload stream, never touching local disk."""
    filename = (file_storage.filename or "").lower()
    try:
        if filename.endswith(".csv"):
            return pd.read_csv(file_storage, low_memory=False)
        elif filename.endswith((".xlsx", ".xls")):
            return pd.read_excel(file_storage, engine="openpyxl")
        return None
    except Exception as e:
        print("Load error:", e)
        return None


def load_csv_bytes(data: bytes):
    return pd.read_csv(io.BytesIO(data), low_memory=False)


# ---------------- COLUMN DETECTION ---------------- #
def detect_columns(df):
    return {
        "numeric": [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])],
        "category": [c for c in df.columns if not pd.api.types.is_numeric_dtype(df[c])]
    }


# ---------------- CLEAN DATA ---------------- #
def clean_data(df, numeric_cols):
    df = df.copy()
    for c in numeric_cols:
        df[c] = pd.to_numeric(df[c], errors="coerce")
    df = df.dropna()
    return df


# ---------------- OUTLIERS ---------------- #
def detect_outliers(df, numeric_cols):
    outliers = {}
    for col in numeric_cols:
        try:
            q1 = df[col].quantile(0.25)
            q3 = df[col].quantile(0.75)
            iqr = q3 - q1
            lower = q1 - 1.5 * iqr
            upper = q3 + 1.5 * iqr
            outliers[col] = int(df[(df[col] < lower) | (df[col] > upper)][col].count())
        except Exception as e:
            print(f"Outlier error in {col}: {e}")
            outliers[col] = 0
    return outliers


# ---------------- KPIs ---------------- #
def generate_kpis(df, numeric_cols):
    kpis = {}
    for col in numeric_cols:
        try:
            kpis[col] = {
                "Total": float(df[col].sum()),
                "Average": float(df[col].mean()),
                "Max": float(df[col].max()),
                "Min": float(df[col].min())
            }
        except Exception:
            continue
    return kpis


# ---------------- INSIGHTS ---------------- #
def generate_text_insights(df, numeric_cols):
    outliers = detect_outliers(df, numeric_cols)
    insights = {}
    for col in numeric_cols:
        try:
            if len(df[col]) < 2:
                insights[col] = "Not enough data for analysis"
                continue
            start = df[col].iloc[0]
            end = df[col].iloc[-1]
            change_text = "N/A" if start == 0 else f"{((end - start) / start) * 100:.2f}%"
            trend = "Increasing" if end > start else "Decreasing"
            insights[col] = (
                f"{col} Analysis:\n\n"
                f"• Average: {df[col].mean():.2f}\n"
                f"• Max: {df[col].max():.2f}\n"
                f"• Min: {df[col].min():.2f}\n"
                f"• Outliers: {outliers.get(col, 0)}\n\n"
                f"• Trend: {trend}\n"
                f"• Change: {change_text}\n\n"
                f"Conclusion:\n"
                f"{col} shows {'strong variation' if df[col].std() > df[col].mean() * 0.5 else 'stable behavior'}."
            )
        except Exception as e:
            print(f"Insight error in {col}: {e}")
            insights[col] = "Error generating insight"
    return insights


# ---------------- BUSINESS SUMMARY ---------------- #
def generate_business_summary(df, numeric_cols):
    if not numeric_cols:
        return "No numeric data available for summary."
    best = max(numeric_cols, key=lambda c: df[c].mean())
    worst = min(numeric_cols, key=lambda c: df[c].mean())
    return (
        f"BUSINESS SUMMARY:\n\n"
        f"• Metrics analyzed: {len(numeric_cols)}\n"
        f"• Best performing metric: {best}\n"
        f"• Weakest metric: {worst}\n\n"
        f"Insight:\nThe dataset shows uneven performance across variables.\n\n"
        f"Recommendation:\nImprove {worst} and scale strategies from {best}."
    )


# ---------------- CORRELATION (in-memory PNG) ---------------- #
def generate_correlation_heatmap_bytes(df):
    if df.shape[1] < 2:
        return None
    corr = df.corr(numeric_only=True)
    if corr.empty or corr.shape[0] < 2:
        return None

    plt.figure(figsize=(6, 5))
    plt.imshow(corr, cmap="coolwarm")
    plt.colorbar()
    plt.xticks(range(len(corr.columns)), corr.columns, rotation=90)
    plt.yticks(range(len(corr.columns)), corr.columns)
    plt.tight_layout()

    buf = io.BytesIO()
    plt.savefig(buf, format="png")
    plt.close()
    buf.seek(0)
    return buf.getvalue()


# ---------------- CHARTS (in-memory PNGs) ---------------- #
def generate_charts_bytes(df, numeric_cols):
    charts = {}
    for col in numeric_cols[:5]:
        try:
            plt.figure()
            df[col].plot(kind="line", title=col)
            buf = io.BytesIO()
            plt.savefig(buf, format="png")
            plt.close()
            buf.seek(0)
            charts[col] = buf.getvalue()
        except Exception as e:
            print(f"Chart error {col}: {e}")
    return charts


# ---------------- THEME ---------------- #
THEME_SETTINGS = {
    "classic": {"bg": (245, 245, 220), "title_color": (101, 67, 33), "text_color": (85, 85, 85), "accent": (222, 184, 135)},
    "modern": {"bg": (102, 126, 234), "title_color": (255, 255, 255), "text_color": (240, 248, 255), "accent": (138, 43, 226)},
    "minimal": {"bg": (224, 231, 255), "title_color": (25, 25, 112), "text_color": (47, 79, 79), "accent": (135, 206, 250)},
    "corporate": {"bg": (15, 23, 42), "title_color": (255, 255, 255), "text_color": (176, 196, 222), "accent": (32, 178, 170)},
    "creative": {"bg": (236, 72, 153), "title_color": (255, 255, 255), "text_color": (255, 228, 225), "accent": (255, 20, 147)},
}


def apply_theme(slide, theme):
    settings = THEME_SETTINGS.get(theme, THEME_SETTINGS["classic"])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*settings["bg"])

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*settings["text_color"])
                    run.font.size = Inches(0.3) if shape == slide.shapes.title else Inches(0.25)

        if shape == slide.shapes.title:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*settings["accent"])
            shape.line.color.rgb = RGBColor(*settings["title_color"])


# ---------------- PPT CREATION (returns bytes, no disk) ---------------- #
def create_presentation_bytes(theme, data, df, chart_bytes_by_col, correlation_bytes):
    if not isinstance(df, pd.DataFrame):
        raise ValueError("df must be a pandas DataFrame")
    if df.empty:
        raise ValueError("DataFrame is empty after cleaning")

    prs = Presentation()
    insights = data.get("insights", {})

    # ---------------- COVER ---------------- #
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Data Analysis Report"
    slide.placeholders[1].text = data.get("source_file", "")
    apply_theme(slide, theme)

    # ---------------- DATA STORY ---------------- #
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Story"
    slide.placeholders[1].text = (
        f"Rows: {len(df)}\nColumns: {len(df.columns)}\n\nIncludes:\n"
        f"• KPIs\n• Trends\n• Outliers\n• Correlation analysis"
    )
    apply_theme(slide, theme)

    # ---------------- KPI SLIDES ---------------- #
    for col, val in data.get("kpis", {}).items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = col
        slide.placeholders[1].text = "\n".join([f"{k}: {v:.2f}" for k, v in val.items()])
        apply_theme(slide, theme)

    # ---------------- INSIGHTS + SIDE CHARTS ---------------- #
    for col, img_bytes in chart_bytes_by_col.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{col} Analysis"

        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
        box.text_frame.text = insights.get(col, "No insight available")

        if img_bytes:
            slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(5.2), Inches(1.5), width=Inches(4.5))

        apply_theme(slide, theme)

    # ---------------- BUSINESS SUMMARY ---------------- #
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Summary"
    slide.placeholders[1].text = data.get("business_summary", "")
    apply_theme(slide, theme)

    # ---------------- CORRELATION ---------------- #
    if correlation_bytes:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Correlation Heatmap"
        slide.shapes.add_picture(io.BytesIO(correlation_bytes), Inches(1), Inches(1.5), width=Inches(8))
        apply_theme(slide, theme)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()


# ---------------- PROCESS FILE ---------------- #
def process_file(file_storage):
    """Parses the upload in-memory, runs analysis, uploads charts + the
    cleaned dataset to blob storage (so /generate_ppt in a *later* request
    can retrieve them), and returns a small JSON-safe summary."""
    df = load_file_from_upload(file_storage)
    if df is None or df.empty:
        return {"status": "error", "message": "Invalid file"}

    detected = detect_columns(df)
    cleaned = clean_data(df, detected["numeric"])
    if cleaned.empty:
        return {"status": "error", "message": "No usable data"}

    report_id = uuid.uuid4().hex
    source_file = file_storage.filename

    kpis = generate_kpis(cleaned, detected["numeric"])
    insights = generate_text_insights(cleaned, detected["numeric"])
    business_summary = generate_business_summary(cleaned, detected["numeric"])

    chart_bytes = generate_charts_bytes(cleaned, detected["numeric"])
    correlation_bytes = generate_correlation_heatmap_bytes(cleaned)

    # Upload chart images so they can be (a) shown on the preview page and
    # (b) re-downloaded during /generate_ppt, without needing local disk.
    chart_urls = {
        col: storage.put_bytes(f"charts/{report_id}/{col}.png", img, "image/png")
        for col, img in chart_bytes.items()
    }
    correlation_url = None
    if correlation_bytes:
        correlation_url = storage.put_bytes(f"charts/{report_id}/correlation.png", correlation_bytes, "image/png")

    # The cleaned dataset itself is needed again in /generate_ppt (for the
    # business summary + slide count) — persist it the same way.
    cleaned_csv_bytes = cleaned.to_csv(index=False).encode("utf-8")
    cleaned_data_url = storage.put_bytes(f"data/{report_id}/cleaned.csv", cleaned_csv_bytes, "text/csv")

    result = {
        "status": "success",
        "report_id": report_id,
        "source_file": source_file,
        "rows": len(cleaned),
        "columns": list(cleaned.columns),
        "kpis": kpis,
        "charts": chart_urls,
        "correlation": correlation_url,
        "insights": insights,
        "business_summary": business_summary,
        "cleaned_data_url": cleaned_data_url,
    }

    # Store the *whole* result as one JSON blob too, and keep only its URL in
    # the session cookie — the session itself has no meaningful size budget
    # for insights + several chart URLs once you have a handful of columns.
    data_blob_url = storage.put_bytes(
        f"session-data/{report_id}.json", json.dumps(result).encode("utf-8"), "application/json"
    )
    result["data_blob_url"] = data_blob_url
    return result


def generate_ppt_from_session_data(theme, data):
    """data is the JSON dict fetched back from data_blob_url. Re-downloads
    the cleaned dataset + chart images from blob storage, builds the deck
    in memory, and returns its bytes."""
    cleaned_bytes = storage.fetch_bytes(data["cleaned_data_url"])
    df = load_csv_bytes(cleaned_bytes)

    chart_bytes_by_col = {
        col: storage.fetch_bytes(url) for col, url in data.get("charts", {}).items()
    }
    correlation_bytes = storage.fetch_bytes(data["correlation"]) if data.get("correlation") else None

    return create_presentation_bytes(theme, data, df, chart_bytes_by_col, correlation_bytes)
